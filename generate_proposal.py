"""
Proposal deck generator for Global HR Staffing Services Pte., Ltd.
Document Management System (DMS) proposal.

USAGE
-----
1. Install the dependency (one time):
       python -m pip install python-pptx
2. Run:
       python generate_proposal.py
3. Output:
       GlobalHR_DMS_Proposal.pptx  (created next to this script)

All text marked with [EDIT ...] is a placeholder you should adjust
(pricing, timeline, your contact details, etc.).
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --------------------------------------------------------------------------
# Branding (colours sampled from the Global HR logo)
# --------------------------------------------------------------------------
BLUE = RGBColor(0x1B, 0x75, 0xBB)      # "GLOBAL" blue
INDIGO = RGBColor(0x2E, 0x31, 0x92)    # "HR" dark indigo
DARK = RGBColor(0x22, 0x29, 0x3B)      # near-black text
GREY = RGBColor(0x5A, 0x5A, 0x5A)      # body grey
LIGHT = RGBColor(0xF2, 0xF6, 0xFB)     # light panel background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Path to the client logo (saved into the Cursor workspace asset folder).
# If the file is missing the title slide simply skips the image.
LOGO_PATH = os.path.join(
    os.path.expanduser("~"),
    ".cursor", "projects", "c-Users-antth-OneDrive-Desktop-Project-DMS",
    "assets",
    "c__Users_antth_AppData_Roaming_Cursor_User_workspaceStorage_"
    "empty-window_images_logo-30c9e653-3535-4f2b-a0f0-5bb94ae67e78.png",
)
# Fallback: a logo.png sitting next to this script.
if not os.path.exists(LOGO_PATH):
    _local = os.path.join(os.path.dirname(os.path.abspath(__file__)), "logo.png")
    if os.path.exists(_local):
        LOGO_PATH = _local

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# --------------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------------
def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def set_run(run, text, size, color, bold=False, italic=False, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def para(tf, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
         space_after=6, bullet=False, level=0, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    prefix = "•  " if bullet else ""
    set_run(p.add_run(), prefix + text, size, color, bold, italic)
    return p


def header(slide, title, kicker=None):
    """Standard content-slide header: accent bar + title."""
    rect(slide, Inches(0), Inches(0), Inches(0.28), SLIDE_H, INDIGO)
    rect(slide, Inches(0.6), Inches(0.55), Inches(0.9), Inches(0.12), BLUE)
    _, tf = textbox(slide, Inches(0.55), Inches(0.65), Inches(12), Inches(1.0))
    para(tf, title, 30, DARK, bold=True, first=True, space_after=0)
    if kicker:
        _, tf2 = textbox(slide, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5))
        para(tf2, kicker, 14, BLUE, italic=True, first=True, space_after=0)


def footer(slide, page):
    _, tf = textbox(slide, Inches(0.55), Inches(7.0), Inches(9), Inches(0.4))
    para(tf, "Global HR Staffing Services Pte., Ltd.  |  DMS Proposal (Confidential)",
         9, GREY, first=True, space_after=0)
    _, tf2 = textbox(slide, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.4))
    para(tf2, str(page), 9, GREY, align=PP_ALIGN.RIGHT, first=True, space_after=0)


def simple_table(slide, x, y, w, rows, col_widths, header_fill=INDIGO,
                 row_h=Inches(0.55), font_size=13):
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_h = row_h * n_rows
    gt = slide.shapes.add_table(n_rows, n_cols, x, y, w, table_h).table
    for ci, cw in enumerate(col_widths):
        gt.columns[ci].width = cw
    for ri, row in enumerate(rows):
        gt.rows[ri].height = row_h
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                set_run(run, str(val), font_size + 1, WHITE, bold=True)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                set_run(run, str(val), font_size, DARK)
    return gt


# --------------------------------------------------------------------------
# Slide 1 — Title
# --------------------------------------------------------------------------
s = add_slide()
rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.35), INDIGO)
rect(s, Inches(0), SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), BLUE)

if os.path.exists(LOGO_PATH):
    s.shapes.add_picture(LOGO_PATH, Inches(4.27), Inches(0.9), width=Inches(4.8))

_, tf = textbox(s, Inches(1), Inches(3.2), Inches(11.33), Inches(1.6),
                anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Document Management System", 40, DARK, bold=True,
     align=PP_ALIGN.CENTER, first=True, space_after=4)
para(tf, "A secure, searchable document platform for the recruitment workflow",
     18, BLUE, align=PP_ALIGN.CENTER, space_after=0)

_, tf = textbox(s, Inches(1), Inches(5.4), Inches(11.33), Inches(1.2),
                anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Proposal prepared for Global HR Staffing Services Pte., Ltd.",
     15, GREY, align=PP_ALIGN.CENTER, first=True, space_after=4)
para(tf, "Prepared by [EDIT: Your Name / Company]   |   [EDIT: Date]",
     13, GREY, align=PP_ALIGN.CENTER, space_after=0)

# --------------------------------------------------------------------------
# Slide 2 — Agenda
# --------------------------------------------------------------------------
s = add_slide()
header(s, "Agenda")
_, tf = textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.6))
items = [
    "The challenge today",
    "The proposed solution",
    "How it works (architecture)",
    "Key features & benefits",
    "Technology & security",
    "Project phases & timeline",
    "Investment",
    "Next steps",
]
for i, it in enumerate(items):
    para(tf, f"{i+1}.   {it}", 18, DARK, first=(i == 0), space_after=10)
footer(s, 2)

# --------------------------------------------------------------------------
# Slide 3 — The Problem
# --------------------------------------------------------------------------
s = add_slide()
header(s, "The Challenge Today", "Recruitment runs on documents — and documents are hard to control")
_, tf = textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.6))
probs = [
    "Thousands of CVs, contracts and compliance documents scattered across folders, drives and email.",
    "Finding the right document quickly is slow and frustrating.",
    "No automatic tracking of expiring compliance documents (e.g. work passes, contracts).",
    "Files can be lost or orphaned when staff leave the company.",
    "No audit trail of who accessed what — a data-protection / compliance risk.",
    "Hard to scale as candidate and client volumes grow.",
]
for i, p in enumerate(probs):
    para(tf, p, 16, DARK, bullet=True, first=(i == 0), space_after=12)
footer(s, 3)

# --------------------------------------------------------------------------
# Slide 4 — The Solution
# --------------------------------------------------------------------------
s = add_slide()
header(s, "The Proposed Solution", "One secure web app for every document")
_, tf = textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.6))
sol = [
    "A custom web application that centralises every document in one secure place.",
    "Powerful search & filtering by candidate, job, client and document type.",
    "Automatic expiry alerts so compliance deadlines are never missed.",
    "Role-based access — admins and recruiters see only what they should.",
    "Full audit log of every upload, view and download.",
    "Files remain browsable directly in company-owned Google Drive as a familiar fallback.",
]
for i, p in enumerate(sol):
    para(tf, p, 16, DARK, bullet=True, first=(i == 0), space_after=12)
footer(s, 4)

# --------------------------------------------------------------------------
# Slide 5 — How It Works (architecture)
# --------------------------------------------------------------------------
s = add_slide()
header(s, "How It Works", "A proven, enterprise-grade Google architecture")

box_y = Inches(2.6)
box_h = Inches(1.7)
box_w = Inches(3.4)
gap = Inches(0.55)
x0 = Inches(0.75)


def arch_box(x, title, lines, fill):
    rect(s, x, box_y, box_w, box_h, fill)
    _, tf = textbox(s, x, box_y + Inches(0.15), box_w, box_h - Inches(0.3),
                    anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, 16, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=6)
    for ln in lines:
        para(tf, ln, 11.5, WHITE, align=PP_ALIGN.CENTER, space_after=2)


arch_box(x0, "Web App",
         ["Browser dashboard", "Login • upload", "search • alerts"], BLUE)
arch_box(x0 + box_w + gap, "Firebase",
         ["Auth + Firestore", "Cloud Functions", "(secure bridge)"], INDIGO)
arch_box(x0 + 2 * (box_w + gap), "Google Drive",
         ["Business Standard", "Shared Drive (2TB+)", "actual PDF / CV files"], BLUE)

# arrows
from pptx.enum.shapes import MSO_SHAPE
for i in range(2):
    ax = x0 + box_w + i * (box_w + gap) + Inches(0.05)
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, box_y + Inches(0.65),
                           gap - Inches(0.1), Inches(0.4))
    a.fill.solid()
    a.fill.fore_color.rgb = GREY
    a.line.fill.background()
    a.shadow.inherit = False

_, tf = textbox(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.6))
para(tf, "Staff log in to a clean web dashboard. Document details live in Firebase "
         "for instant search; the files themselves are stored safely in your "
         "company-owned Google Workspace Drive. Everything is secure, tracked and "
         "scalable — and staff can still open the Drive folders directly.",
     14, GREY, first=True, space_after=0)
footer(s, 5)

# --------------------------------------------------------------------------
# Slide 6 — Key Features
# --------------------------------------------------------------------------
s = add_slide()
header(s, "Key Features & Benefits")
rows = [
    ["Feature", "Benefit to Global HR"],
    ["Secure staff login", "Only authorised people access candidate data"],
    ["Smart search & filters", "Find any document in seconds"],
    ["Expiry alerts", "Never miss a compliance / work-pass deadline"],
    ["Candidate profiles", "Every document for a person in one view"],
    ["Audit log", "Full compliance & data-protection trail"],
    ["Google Drive browsing", "Familiar fallback access for staff"],
]
simple_table(s, Inches(0.9), Inches(2.0), Inches(11.5), rows,
             [Inches(4.0), Inches(7.5)], row_h=Inches(0.62))
footer(s, 6)

# --------------------------------------------------------------------------
# Slide 7 — Technology & Security
# --------------------------------------------------------------------------
s = add_slide()
header(s, "Technology & Security", "Built on infrastructure trusted by global enterprises")
_, tf = textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.6))
tech = [
    "Firebase (by Google) — secure, scalable backend for login and data.",
    "Google Workspace Drive — business-grade file storage (2TB+, pooled).",
    "Modern web app — works in any browser, nothing to install.",
    "Role-based permissions & security rules protect sensitive candidate data.",
    "Hosted on Google Cloud — reliable, backed-up and continuously available.",
    "Designed with data-protection (PDPA / GDPR-style) compliance in mind.",
]
for i, p in enumerate(tech):
    para(tf, p, 16, DARK, bullet=True, first=(i == 0), space_after=12)
footer(s, 7)

# --------------------------------------------------------------------------
# Slide 8 — Timeline
# --------------------------------------------------------------------------
s = add_slide()
header(s, "Project Phases & Timeline", "[EDIT: adjust durations to suit your schedule]")
rows = [
    ["Phase", "Deliverable", "Estimate"],
    ["1. Setup & Design", "Firebase project, data model, UI design", "[EDIT] wks"],
    ["2. Core Application", "Login, candidate records, upload", "[EDIT] wks"],
    ["3. Drive Integration", "Secure file storage bridge", "[EDIT] wks"],
    ["4. Search & Alerts", "Filters, expiry tracking, dashboard", "[EDIT] wks"],
    ["5. Testing & Launch", "Testing, deployment, staff training", "[EDIT] wks"],
]
simple_table(s, Inches(0.9), Inches(2.1), Inches(11.5), rows,
             [Inches(3.0), Inches(6.3), Inches(2.2)], row_h=Inches(0.62))
footer(s, 8)

# --------------------------------------------------------------------------
# Slide 9 — Investment
# --------------------------------------------------------------------------
s = add_slide()
header(s, "Investment", "All amounts in SGD — [EDIT placeholder figures]")
rows = [
    ["Item", "Type", "Amount (SGD)"],
    ["DMS design & development", "One-time", "$[EDIT]"],
    ["Google Workspace Business Standard", "Per user / month", "~$28*"],
    ["Firebase (hosting + functions)", "Monthly", "~$[EDIT] (low)"],
    ["Support & maintenance (optional)", "Monthly", "$[EDIT]"],
]
simple_table(s, Inches(0.9), Inches(2.1), Inches(11.5), rows,
             [Inches(5.5), Inches(3.2), Inches(2.8)], row_h=Inches(0.62))
_, tf = textbox(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.2))
para(tf, "* Google Workspace is billed by Google per user. Storage is pooled across "
         "users, comfortably covering 1TB+ of documents. Firebase has a generous free "
         "tier; running costs stay low at typical agency volumes.",
     12, GREY, italic=True, first=True, space_after=0)
footer(s, 9)

# --------------------------------------------------------------------------
# Slide 10 — Why it pays off
# --------------------------------------------------------------------------
s = add_slide()
header(s, "Why This Investment Pays Off")
_, tf = textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.6))
pay = [
    "Saves hours of manual searching every week across the team.",
    "Avoids costly compliance breaches, fines and reputational damage.",
    "Protects the business against lost or orphaned documents.",
    "Scales smoothly as the agency takes on more candidates and clients.",
    "A professional, modern system that clients and candidates notice.",
]
for i, p in enumerate(pay):
    para(tf, p, 17, DARK, bullet=True, first=(i == 0), space_after=14)
footer(s, 10)

# --------------------------------------------------------------------------
# Slide 11 — Next Steps
# --------------------------------------------------------------------------
s = add_slide()
header(s, "Next Steps")
_, tf = textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.6))
steps = [
    "Approve this proposal and confirm budget.",
    "Set up Google Workspace and the Firebase project.",
    "Begin Phase 1 — setup & design.",
    "Agree review checkpoints for each phase.",
]
for i, st in enumerate(steps):
    para(tf, f"{i+1}.   {st}", 18, DARK, first=(i == 0), space_after=14)
footer(s, 11)

# --------------------------------------------------------------------------
# Slide 12 — Thank you
# --------------------------------------------------------------------------
s = add_slide()
rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, INDIGO)
rect(s, Inches(0), SLIDE_H - Inches(0.4), SLIDE_W, Inches(0.4), BLUE)
_, tf = textbox(s, Inches(1), Inches(2.6), Inches(11.33), Inches(2.4),
                anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Thank You", 48, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True,
     space_after=10)
para(tf, "Let's build a secure, scalable document system for Global HR.",
     18, WHITE, align=PP_ALIGN.CENTER, space_after=20)
para(tf, "[EDIT: Your Name]  •  [EDIT: email]  •  [EDIT: phone]",
     15, RGBColor(0xCF, 0xD8, 0xF0), align=PP_ALIGN.CENTER, space_after=0)

# --------------------------------------------------------------------------
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "GlobalHR_DMS_Proposal.pptx")
prs.save(out)
print("Created:", out)
